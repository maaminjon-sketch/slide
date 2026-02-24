import streamlit as st
import os
import json
import io
import re
import base64
from urllib.parse import quote
from groq import Groq
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from dotenv import load_dotenv
import gspread
from oauth2client.service_account import ServiceAccountCredentials
from datetime import datetime
import google.generativeai as genai
from google.generativeai.types import HarmCategory, HarmBlockThreshold
import requests

WORD_MAP = {
    20: 120,
    21: 110,
    22: 110,
    23: 101,
    24: 95,
    25: 89,
    26: 82,
    27: 77,
    28: 71,
    29: 60,
    30: 55
}

STOPWORDS = {
    "the", "and", "for", "with", "from", "into", "about", "this", "that", "these", "those", "are", "was",
    "were", "have", "has", "had", "its", "their", "your", "our", "как", "что", "это", "эти", "для", "при",
    "или", "also", "than", "then", "when", "where", "which", "who", "whom", "whose", "как", "про", "об",
    "над", "под", "без", "after", "before", "during", "роль", "история", "введение", "заключение"
}

APP_BRAND = "SLIDEX-369"
APP_CLOUD_URL = "https://tinyurl.com/amin-cloud"
APP_FOUNDER = "Shodmehr vs Amin"
APP_ICON_PATH = "icon.png"


# Загрузка ключей из .env
load_dotenv()
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
client = Groq(api_key=os.getenv("GROQ_API_KEY"))

# --- ВСПОМОГАТЕЛЬНАЯ ФУНКЦИЯ ДЛЯ JSON ---
def extract_json(text):
    """Находит и возвращает валидный JSON объект из строки, игнорируя markdown."""
    try:
        # Ищем первую открывающую { и последнюю закрывающую }
        start = text.find('{')
        end = text.rfind('}') + 1
        
        if start != -1 and end != 0:
            json_str = text[start:end]
            return json.loads(json_str)
        return None
    except json.JSONDecodeError:
        return None

def verify_and_update_key(user_key, mode="check"):
    try:
        scope = ["https://spreadsheets.google.com/feeds", "https://www.googleapis.com/auth/drive"]
        
        # Берем данные из st.secrets (TOML), а не из файла credentials.json
        creds_info = st.secrets["gcp_service_account"]
        creds = ServiceAccountCredentials.from_json_keyfile_dict(dict(creds_info), scope)
        client_gs = gspread.authorize(creds)
        
        sheet = client_gs.open("SlideGen_DB").sheet1
        keys_list = sheet.col_values(1) 
        
        if user_key not in keys_list:
            return False, "Ключ не найден!"
            
        row_idx = keys_list.index(user_key) + 1
        row_data = sheet.row_values(row_idx)
        
        expiry_date = datetime.strptime(row_data[1], "%Y-%m-%d")
        if datetime.now() > expiry_date:
            return False, "Срок ключа истек!"
            
        limit = int(row_data[2])
        if limit <= 0:
            return False, "Лимиты исчерпаны!"
            
        if mode == "update":
            new_limit = limit - 1
            sheet.update_cell(row_idx, 3, new_limit) 
            return True, new_limit
        
        return True, limit
    except Exception as e:
        return False, f"Ошибка базы: {str(e)}"

# --- ВСПОМОГАТЕЛЬНЫЕ ФУНКЦИИ ---
def extract_json(text):
    try:
        start = text.find('{')
        end = text.rfind('}') + 1
        if start != -1 and end != 0:
            return json.loads(text[start:end])
        return None
    except: return None

# --- ФУНКЦИЯ ГЕНЕРАЦИИ ТЕКСТА ЧЕРЕЗ GROQ ---
def get_ai_content(topic, lang, slides_count, user_font_size):
    target_words = WORD_MAP.get(user_font_size, 60)
    prompt = f"""
    You are a senior presentation writer.
    Create a professional and academic presentation about "{topic}" in {lang}. 
    Total slides: {slides_count}.

    STRICT OUTPUT: Return ONLY valid JSON object. No markdown, no comments.

    REQUIRED JSON FORMAT:
    {{
      "slides": [
        {{
          "title": "Slide title",
          "content": ["Bullet sentence 1", "Bullet sentence 2", "Bullet sentence 3"],
          "image_query": "Precise English photo query"
        }}
      ]
    }}

    IMPORTANT RULES:
    1. PROGRESSION: Every slide must cover a DIFFERENT subtopic in logical sequence.
    2. NO REPETITION: Never repeat facts, phrases, or ideas between slides.
    3. DEPTH: 6-8 informative bullet sentences per slide.
    4. LENGTH: Each slide must contain approximately {target_words} words total in "content".
    5. QUALITY: Each bullet should be concrete, factual, and useful. Avoid filler phrases.
    6. IMAGE QUERY: For each slide provide a specific search query focused on main entity/action,
       not generic geography. Example: if slide is about Churchill's role in England,
       query should target "Winston Churchill portrait WWII speech", not "England".
    7. Slide 1 = introduction context. Last slide = conclusion with takeaways.

    Return ONLY a JSON object.
    """
    
    try:
        completion = client.chat.completions.create(
            model="llama-3.1-8b-instant", 
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7,
            max_tokens=4000, 
            response_format={ "type": "json_object" }
        )
        return completion.choices[0].message.content
    except Exception as e:
        return str(e)

# --- GEMINI (ИСПРАВЛЕНО: НОВАЯ ТАБЛИЦА СЛОВ) ---
def get_gemini_content(topic, lang, slides_count, user_font_size):
    model = genai.GenerativeModel('gemini-2.0-flash') 
    
    # Если шрифт не в списке, берем дефолт 60
    target_words = WORD_MAP.get(user_font_size, 60)

    prompt = f"""
    You are a professional presentation generator. Create a structured presentation about "{topic}" in {lang}.
    Total slides: {slides_count}.

    STRICT JSON OUTPUT FORMAT ONLY. NO MARKDOWN. NO ```json wrappers.
    
    Structure:
    {{
      "slides": [
        {{
          "title": "Clear Title Here",
          "content": ["Sentence 1...", "Sentence 2...", "Sentence 3..."],
          "image_query": "description for image search or null"
        }}
      ]
    }}

    CONTENT RULES:
    1. Each slide must contain {target_words} to {target_words + 8} words TOTAL across all bullet points. THIS IS CRITICAL.
    2. Slide 1 is Introduction. Last slide is Conclusion.
    3. Use factual, academic tone.
    4. "content" MUST be an array of strings (bullet points).
    5. "image_query": Provide a short English query for every slide (never null).
    6. NO REPETITION: no repeated facts or repeated wording across slides.
    7. Build image_query around PRIMARY ENTITY of the slide.
       Example: if text discusses Churchill and leadership in wartime Britain,
       use "Winston Churchill WWII portrait speech" and do NOT use generic "England landscape".
    8. Add concrete facts in each slide: named events, people, organizations, metrics, and real-world details.
    9. Include years/dates and numeric evidence where relevant (history, economics, politics, science, technology, etc.).
    10. If topic is quantitative (history, tajik,math, physics, engineering, finance, chemistry, statistics), include formulas or equation-like expressions where useful.
    11. Avoid vague text. Every bullet must provide specific information, not generic statements.
    12. Minimum 5 bullet points per slide, each bullet should contain specific factual detail.
    13. For history/politics slides, include at least one named person and one dated event per slide where possible.
    14. image_query MUST include the exact proper name for people (e.g., "Gamal Abdel Nasser portrait 1956") when a person is central.
    """

    # Настройки безопасности (отключаем блокировку)
    safety_settings = {
        HarmCategory.HARM_CATEGORY_HATE_SPEECH: HarmBlockThreshold.BLOCK_NONE,
        HarmCategory.HARM_CATEGORY_HARASSMENT: HarmBlockThreshold.BLOCK_NONE,
        HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT: HarmBlockThreshold.BLOCK_NONE,
        HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT: HarmBlockThreshold.BLOCK_NONE,
    }

    try:
        response = model.generate_content(
            prompt,
            safety_settings=safety_settings,
            generation_config=genai.types.GenerationConfig(
                response_mime_type="application/json",
                temperature=0.3,
                max_output_tokens=8000
            )
        )
        return response.text
    except Exception as e:
        print(f"Gemini Error: {e}")
        return "{}"


def clean_text(text):
    if not text: return ""
    # Убираем <B>, **, __ и прочий мусор
    return re.sub(r'<[^>]+>|\*\*|__', '', str(text)).strip()

def load_logo_asset(path):
    """Загружает лого и обрезает прозрачные поля, чтобы значок не выглядел маленьким."""
    if not path or not os.path.exists(path):
        return None
    try:
        from PIL import Image
        img = Image.open(path).convert("RGBA")
        alpha = img.split()[-1]
        bbox = alpha.getbbox()
        if bbox:
            img = img.crop(bbox)
        return img
    except Exception:
        # Fallback: Streamlit сам загрузит файл, даже если PIL недоступен
        return path

def logo_data_uri(logo_asset):
    """Преобразует лого в data URI для корректного HTML-рендера без ограничений Streamlit image."""
    if logo_asset is None:
        return None
    try:
        if isinstance(logo_asset, str):
            with open(logo_asset, "rb") as f:
                b64 = base64.b64encode(f.read()).decode("utf-8")
            return f"data:image/png;base64,{b64}"

        # PIL Image
        buf = io.BytesIO()
        logo_asset.save(buf, format="PNG")
        b64 = base64.b64encode(buf.getvalue()).decode("utf-8")
        return f"data:image/png;base64,{b64}"
    except Exception:
        return None

def get_effective_slide_query(slide_data, topic):
    model_query = clean_text(slide_data.get("image_query", ""))
    raw_title = clean_text(slide_data.get("title", ""))
    raw_content = slide_data.get("content", [])
    if isinstance(raw_content, list):
        raw_content = " ".join([clean_text(x) for x in raw_content if x])
    else:
        raw_content = clean_text(raw_content)

    source = " ".join([model_query, raw_title, raw_content, clean_text(topic)])
    tokens = re.findall(r"[A-Za-zА-Яа-я0-9][A-Za-zА-Яа-я0-9-]{2,}", source)
    filtered = []
    for token in tokens:
        low = token.lower()
        if low in STOPWORDS:
            continue
        if len(low) < 3:
            continue
        if low not in filtered:
            filtered.append(low)

    # Если запрос уже выглядит как запрос про персону, не засоряем его хвостом.
    person_like = bool(re.search(r"\b[A-Z][a-z]+\s+[A-Z][a-z]+\b", model_query))

    # Приоритет: берем query от модели + ключевые токены из контекста
    if model_query:
        model_base = model_query.strip()
    else:
        model_base = raw_title if raw_title else clean_text(topic)

    if person_like:
        return re.sub(r"\s+", " ", model_base)

    tail = " ".join(filtered[:5])
    final_query = (f"{model_base} {tail}").strip()
    return re.sub(r"\s+", " ", final_query)

def score_duckduckgo_result(result, query_tokens):
    score = 0
    title = str(result.get("title", "")).lower()
    source = str(result.get("source", "")).lower()
    page_url = str(result.get("url", "")).lower()
    image_url = str(result.get("image", "")).lower()
    haystack = f"{title} {source} {page_url} {image_url}"

    for token in query_tokens:
        if token in haystack:
            score += 2

    # Предпочитаем картинки с нормальным соотношением сторон для слайда.
    try:
        w = int(result.get("width") or 0)
        h = int(result.get("height") or 0)
        if w >= 1000 and h >= 550:
            score += 2
        if h > 0:
            ratio = w / h
            if 1.2 <= ratio <= 2.2:
                score += 2
    except:
        pass

    if image_url.endswith((".jpg", ".jpeg", ".png")):
        score += 1
    if image_url.endswith((".svg", ".webp")):
        score -= 3

    return score

def get_duckduckgo_vqd(query):
    if not query:
        return None
    headers = {
        "User-Agent": "Mozilla/5.0",
        "Accept-Language": "en-US,en;q=0.9"
    }
    try:
        r = requests.get(
            "https://duckduckgo.com/",
            params={"q": query, "iax": "images", "ia": "images"},
            headers=headers,
            timeout=6
        )
        if r.status_code != 200:
            return None
        text = r.text
        patterns = [
            r"vqd='([^']+)'",
            r'vqd="([^"]+)"',
            r'"vqd":"([^"]+)"',
            r"vqd=([0-9-]+)\&"
        ]
        for pattern in patterns:
            match = re.search(pattern, text)
            if match:
                return match.group(1)
    except:
        pass
    return None

def search_duckduckgo_image(query):
    if not query:
        return None

    vqd = get_duckduckgo_vqd(query)
    if not vqd:
        return None

    headers = {
        "User-Agent": "Mozilla/5.0",
        "Accept": "application/json, text/javascript, */*; q=0.01",
        "Referer": f"https://duckduckgo.com/?q={quote(query)}&iax=images&ia=images"
    }
    params = {
        "l": "wt-wt",
        "o": "json",
        "q": query,
        "vqd": vqd,
        "f": ",,,"
    }

    try:
        r = requests.get("https://duckduckgo.com/i.js", headers=headers, params=params, timeout=8)
        if r.status_code != 200:
            return None
        items = r.json().get("results", [])
        if not items:
            return None

        query_tokens = [
            t.lower()
            for t in re.findall(r"[A-Za-zА-Яа-я0-9-]{3,}", query)
            if t.lower() not in STOPWORDS
        ]
        ranked = sorted(items, key=lambda x: score_duckduckgo_result(x, query_tokens), reverse=True)

        for item in ranked:
            image_url = item.get("image") or item.get("thumbnail")
            if not image_url:
                continue
            if image_url.startswith("//"):
                image_url = f"https:{image_url}"
            if not image_url.startswith(("http://", "https://")):
                continue
            low = image_url.lower()
            if low.endswith(".svg"):
                continue
            return image_url
    except:
        pass
    return None

def fetch_ppt_compatible_image(image_url, timeout=7):
    if not image_url:
        return None
    headers = {
        "User-Agent": "Mozilla/5.0",
        "Accept": "image/*,*/*;q=0.8",
        "Referer": "https://duckduckgo.com/"
    }
    try:
        r = requests.get(image_url, headers=headers, timeout=timeout)
        if r.status_code != 200 or not r.content:
            return None
        data = r.content
        ctype = str(r.headers.get("Content-Type", "")).lower()

        # JPEG / PNG напрямую поддерживаются python-pptx.
        if data[:2] == b"\xff\xd8" or data[:8] == b"\x89PNG\r\n\x1a\n":
            return data
        if "jpeg" in ctype or "jpg" in ctype or "png" in ctype:
            return data

        # Пробуем конвертировать WEBP и другие форматы, если Pillow доступен.
        if "svg" in ctype or image_url.lower().endswith(".svg"):
            return None
        try:
            from PIL import Image
            with Image.open(io.BytesIO(data)) as img:
                out = io.BytesIO()
                if img.mode in ("RGBA", "LA", "P"):
                    img.convert("RGBA").save(out, format="PNG")
                else:
                    img.convert("RGB").save(out, format="JPEG")
                return out.getvalue()
        except Exception:
            return None
    except:
        pass
    return None

def search_wikimedia_image(query):
    if not query:
        return None
    try:
        # 1) Ищем страницу в Wikipedia по запросу
        s_url = f"https://en.wikipedia.org/w/api.php?action=query&list=search&srsearch={quote(query)}&utf8=1&format=json"
        r1 = requests.get(s_url, timeout=6)
        if r1.status_code != 200:
            return None
        s_data = r1.json()
        search_items = s_data.get("query", {}).get("search", [])
        if not search_items:
            return None
        title = search_items[0].get("title")
        if not title:
            return None

        # 2) Берем превью картинки страницы (обычно портрет/релевантное фото)
        p_url = (
            "https://en.wikipedia.org/w/api.php?action=query&prop=pageimages"
            f"&piprop=original|thumbnail&pithumbsize=1600&titles={quote(title)}&format=json"
        )
        r2 = requests.get(p_url, timeout=6)
        if r2.status_code != 200:
            return None
        p_data = r2.json()
        pages = p_data.get("query", {}).get("pages", {})
        for page in pages.values():
            original = page.get("original", {})
            thumb = page.get("thumbnail", {})
            if original.get("source"):
                return original["source"]
            if thumb.get("source"):
                return thumb["source"]
    except:
        pass
    return None

def has_person_name_hint(query):
    if not query:
        return False
    # 2+ слов с заглавных букв обычно указывают на имя/фамилию
    latin = bool(re.search(r"\b[A-Z][a-z]+\s+[A-Z][a-z]+\b", query))
    cyrillic = bool(re.search(r"\b[А-ЯЁ][а-яё]+\s+[А-ЯЁ][а-яё]+\b", query))
    return latin or cyrillic

def resolve_slide_image(slide_data, topic):
    query = get_effective_slide_query(slide_data, topic)
    if not query:
        return None

    person_context = has_person_name_hint(query) or any(
        k in query.lower() for k in ["portrait", "leader", "president", "king", "prime minister"]
    )

    candidates = []

    # Для персон и исторических лидеров сначала Wikipedia/Wikimedia (точнее портреты).
    if person_context:
        candidates.append(search_wikimedia_image(query))

    # Затем DuckDuckGo Images.
    candidates.append(search_duckduckgo_image(query))

    # Последний fallback: снова Wikipedia/Wikimedia.
    candidates.append(search_wikimedia_image(query))

    checked = set()
    for image_url in candidates:
        if not image_url or image_url in checked:
            continue
        checked.add(image_url)
        img_data = fetch_ppt_compatible_image(image_url)
        if img_data:
            return img_data
    return None


def create_pptx(ai_json_text, template_folder, user_font_size=20, topic="", include_images=True, image_slide_word_ratio=1.0):
    # Золотое правило (лимит слов)
    limit = WORD_MAP.get(user_font_size, 60)

    style_colors = {
        "yellow_style": (0, 0, 139),
        "modern_red": (255, 255, 255),
        "tech_blue": (0, 255, 255),
        "minimal_white": (35, 35, 35),
        "dark_mode": (200, 200, 200)
    }

    data = extract_json(ai_json_text)
    if not data or "slides" not in data:
        print("Ошибка: Неверный формат JSON или пустой ответ")
        return None

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)

    bg_images = []
    if template_folder and os.path.isdir(template_folder):
        bg_images = sorted([os.path.join(template_folder, f) for f in os.listdir(template_folder) 
                      if f.lower().endswith(('.png', '.jpg', '.jpeg'))])

    folder_name = os.path.basename(os.path.normpath(template_folder)).lower()

    for i, slide_data in enumerate(data["slides"]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Проверка, нужен ли слот для фото
        show_image_space = include_images and (i % 2 == 1)

        # Инициализация image_data всегда, чтобы не было ошибки
        image_data = None
        if show_image_space:
            image_data = resolve_slide_image(slide_data, topic)

        
        # Фон слайда
        if bg_images:
            current_bg = bg_images[i % len(bg_images)]
            slide.shapes.add_picture(current_bg, 0, 0, prs.slide_width, prs.slide_height)
            if folder_name in style_colors:
                r, g, b = style_colors[folder_name]
                text_color = RGBColor(r, g, b)
            else:
                is_dark = "white" not in current_bg.lower() and "light" not in current_bg.lower()
                text_color = RGBColor(255, 255, 255) if is_dark else RGBColor(35, 35, 35)
        else:
            text_color = RGBColor(0, 0, 0)

        accent_color = RGBColor(0, 255, 127)
        b_width = Inches(7.5) if show_image_space else Inches(11.5)

        # Добавление картинки
        if image_data:
            try:
                slide.shapes.add_picture(io.BytesIO(image_data),
                                         Inches(8.6), Inches(1.2),
                                         Inches(4.2), Inches(4.8))
            except:
                pass

        # Заголовок
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(11.5), Inches(1.0))
        p_t = title_box.text_frame.paragraphs[0]
        p_t.text = clean_text(slide_data.get("title", "No Title")).upper()
        p_t.font.size = Pt(32)
        p_t.font.name = "Times New Roman"
        p_t.font.color.rgb = accent_color
        p_t.font.bold = True

        # Контент
        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), b_width, Inches(6.1))
        tf_body = body_box.text_frame
        tf_body.word_wrap = True
        tf_body.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf_body.margin_left = Inches(0.1)
        tf_body.margin_right = Inches(0.1)
        tf_body.margin_top = Inches(0.08)
        tf_body.margin_bottom = Inches(0.08)

        current_word_count = 0
        slide_limit = limit
        if show_image_space and image_slide_word_ratio < 1.0:
            slide_limit = max(20, int(limit * image_slide_word_ratio))

        raw_content = slide_data.get("content", [])
        if isinstance(raw_content, str):
            raw_content = [raw_content]

        cleaned_points = []
        for point in raw_content:
            p_text = clean_text(point)
            if not p_text:
                continue
            cleaned_points.append(p_text)

        # Строго удерживаем общий объем текста в пределах лимита для выбранного шрифта
        fitted_points = []
        for p_text in cleaned_points:
            words = p_text.split()
            if current_word_count + len(words) <= slide_limit:
                fitted_points.append(p_text)
                current_word_count += len(words)
                continue
            remaining = slide_limit - current_word_count
            if remaining >= 8:
                fitted_points.append(" ".join(words[:remaining]).rstrip(",.;:") + ".")
                break
            if remaining < 8:
                break

        if not fitted_points:
            fitted_points = cleaned_points[:1]

        # Группируем по 2 предложения в абзац (не слишком дробно и не одним блоком)
        grouped_points = []
        pair = []
        for t in fitted_points:
            pair.append(t)
            if len(pair) == 2:
                grouped_points.append(" ".join(pair))
                pair = []
        if pair:
            grouped_points.append(" ".join(pair))
        if grouped_points:
            fitted_points = grouped_points

        # Используем первый абзац, чтобы не было пустой строки сверху
        tf_body.clear()
        first = True
        for p_text in fitted_points:
            p = tf_body.paragraphs[0] if first else tf_body.add_paragraph()
            p.text = p_text
            p.font.size = Pt(user_font_size)
            p.font.name = "Times New Roman"
            p.font.color.rgb = text_color
            p.space_after = Pt(2)
            p.line_spacing = 1.12
            first = False

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out


st.markdown("""
    <style>
    /* Главный фон: Темно-серый с глубоким зеленым оттенком */
    .stApp {
        background: radial-gradient(circle at top, #0a1f1a 0%, #020504 100%);
        color: #e0e0e0;
    }

    /* Делаем превью выбора стиля очень маленькими */
    [data-testid="stHorizontalBlock"] [data-testid="stImage"] img {
        border-radius: 10px;
        max-height: 60px; /* Ограничиваем высоту */
        object-fit: contain;
    }

    /* Кнопки стилей */
    [data-testid="stHorizontalBlock"] button {
        padding: 2px 5px !important;
        font-size: 10px !important;
        min-height: 25px !important;
        white-space: nowrap !important;
        overflow: hidden !important;
        text-overflow: ellipsis;
    }

    /* Заголовки */
    h1, h2, h3 {
        color: #00ff7f !important;
        text-shadow: 0 0 10px rgba(0, 255, 127, 0.5);
    }

    /* Бренд вверху: красный стиль + эффект "прыгучести" */
    .slidex-brand {
        text-align: center;
        margin-top: -8px;
        margin-bottom: 8px;
        font-size: 3rem;
        font-weight: 900;
        letter-spacing: 1px;
        color: #ff2a2a !important;
        text-shadow: 0 0 6px rgba(255, 42, 42, 0.8), 0 0 16px rgba(255, 0, 0, 0.45);
        animation: slidex-bounce 2.2s ease-in-out infinite;
    }

    .site-logo-wrap {
        display: flex;
        justify-content: center;
        margin-top: 6px;
        margin-bottom: 6px;
    }

    .site-logo {
        width: 280px;
        max-width: 52vw;
        height: auto;
        object-fit: contain;
        filter: drop-shadow(0 2px 6px rgba(0, 0, 0, 0.45));
    }

    @keyframes slidex-bounce {
        0%, 100% {
            transform: translateY(0);
        }
        20% {
            transform: translateY(-7px);
        }
        45% {
            transform: translateY(0);
        }
        65% {
            transform: translateY(-3px);
        }
    }

    /* Мобильная адаптация */
    @media (max-width: 768px) {
        .site-logo {
            width: 210px;
            max-width: 64vw;
        }

        .slidex-brand {
            font-size: 2.15rem;
            letter-spacing: 0.5px;
            margin-top: 0;
            margin-bottom: 4px;
        }

        div.stButton > button {
            font-size: 1rem !important;
            padding: 0.6em 1em !important;
            border-radius: 10px !important;
        }

        [data-testid="stHorizontalBlock"] {
            gap: 0.45rem !important;
        }

        [data-testid="stHorizontalBlock"] [data-testid="column"] {
            min-width: calc(50% - 0.45rem) !important;
            flex: 1 1 calc(50% - 0.45rem) !important;
        }

        [data-testid="stHorizontalBlock"] [data-testid="stImage"] img {
            max-height: 90px !important;
            border-radius: 8px;
            object-fit: contain !important;
        }

        [data-testid="stHorizontalBlock"] button {
            font-size: 11px !important;
            min-height: 34px !important;
            padding: 6px 8px !important;
        }
    }

    /* КРУТАЯ КНОПКА С ЭФФЕКТАМИ */
    div.stButton > button {
        background: linear-gradient(45deg, #006400, #00ff7f);
        color: white !important;
        border: none !important;
        border-radius: 12px !important;
        padding: 0.75em 2em !important;
        font-size: 1.2rem !important;
        font-weight: bold !important;
        text-transform: uppercase;
        box-shadow: 0 0 15px rgba(0, 255, 127, 0.3);
    }
    div.stButton > button:hover {
        transform: scale(1.05);
        box-shadow: 0 0 25px rgba(0, 255, 127, 0.7);
        color: #ffffff !important;
    }
    
    /* Стилизация Selectbox */
    div[data-baseweb="select"] > div {
        background-color: #0a1f1a !important;
        color: white !important;
    }
    </style>
    """, unsafe_allow_html=True)

logo_asset = load_logo_asset(APP_ICON_PATH)
if logo_asset:
    logo_uri = logo_data_uri(logo_asset)
    if logo_uri:
        st.markdown(
            f"<div class='site-logo-wrap'><img src='{logo_uri}' class='site-logo' alt='SLIDEX logo'></div>",
            unsafe_allow_html=True
        )
    else:
        col_l, col_c, col_r = st.columns([4, 3, 4])
        with col_c:
            st.image(logo_asset, width=220)

st.markdown(
    f"<h1 class='slidex-brand'>{APP_BRAND}</h1>",
    unsafe_allow_html=True
)


# --- Боковая панель ---
st.sidebar.title("💎 Выбор версии")
if os.path.exists(APP_ICON_PATH):
    st.sidebar.image(logo_asset if logo_asset else APP_ICON_PATH, width=120)
mode = st.sidebar.radio("Доступные варианты:", ["Бесплатная (Groq)🆓", "Платная версия PRO 🆕"])
st.sidebar.markdown("---")
st.sidebar.info("Версия 1.0.5 \nПоддержка: [WhatsApp Админа](https://wa.me/992001255656)")
st.sidebar.markdown(f"☁️ [Amin Cloud]({APP_CLOUD_URL})")

# --- ЛОГИКА БЕСПЛАТНОЙ ВЕРСИИ ---
if mode == "Бесплатная (Groq)🆓":
    st.title("🎁 Бесплатная генерация")
    st.write("Создавайте презентации быстро с помощью Groq AI.")
    
    topic = st.text_input("Тема презентации", placeholder="Например: История Таджикистана")
    
    col1, col2 = st.columns(2)
    with col1:
        lang = st.selectbox("Язык", ["Русский", "Таджикский", "English"], key="free_lang")
        
    with col2:
        subdivisions = st.number_input("Кол-во слайдов", 5, 7, 5, key="free_slides")
    
    user_font_size = st.slider("Размер шрифта текста", 20, 30, 22, key="free_font")

    st.write("---")
    st.subheader("🎨 Выберите стиль оформления")

    templates_dir = "templates"
    if not os.path.exists(templates_dir):
        os.makedirs(templates_dir)

    style_folders = [f for f in os.listdir(templates_dir) if os.path.isdir(os.path.join(templates_dir, f))]

    if style_folders:
        f_cols = st.columns(3) 
        for idx, folder_name in enumerate(style_folders):
            col_idx = idx % 3
            with f_cols[col_idx]:
                folder_path = os.path.join(templates_dir, folder_name)
                try:
                    available_imgs = [f for f in os.listdir(folder_path) if f.endswith(('.png', '.jpg', '.jpeg'))]
                    if available_imgs:
                        st.image(os.path.join(folder_path, available_imgs[0]), use_container_width=True)
                except:
                    pass

                if st.button(folder_name.upper(), key=f"style_btn_{idx}", use_container_width=True):
                    st.session_state['selected_tpl_folder'] = folder_path
                    st.session_state['style_name_display'] = folder_name.upper()

        if 'style_name_display' in st.session_state:
            st.markdown(f"<p style='color: #00ff7f;'>✅ Выбран стиль: <b>{st.session_state['style_name_display']}</b></p>", unsafe_allow_html=True)
    else:
        st.warning("В папке 'templates' нет подпапок со стилями.")

    final_template_path = st.session_state.get('selected_tpl_folder', None)
    st.write("---")

    if st.button("Сгенерировать презентацию", key="free_gen_btn"):
            if topic:
                if not os.getenv("GROQ_API_KEY"):
                    st.error("Ошибка: Не найден API ключ Groq в файле .env")
                else:
                    if not final_template_path:
                        st.warning("Пожалуйста, выберите стиль оформления!")
                    else:
                        with st.spinner("🤖 Groq генерирует контент..."):
                            try:
                                ai_json_content = get_ai_content(topic, lang, subdivisions, user_font_size)
                                ppt_file = create_pptx(
                                    ai_json_content,
                                    template_folder=final_template_path,
                                    user_font_size=user_font_size,
                                    topic=topic,
                                    include_images=False,
                                    image_slide_word_ratio=1.0
                                )
                                
                                if ppt_file:
                                    st.success("✅ Презентация готова!")
                                    st.download_button(
                                        label="📥 Скачать презентацию (.pptx)",
                                        data=ppt_file,
                                        file_name=f"{topic}.pptx",
                                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                    )
                                else:
                                    st.error("Ошибка генерации файла. Попробуйте еще раз.")
                            except Exception as e:
                                st.error(f"Произошла ошибка: {e}")
            else:
                st.error("Введите тему презентации!")


# --- ЛОГИКА ПЛАТНОЙ ВЕРСИИ ---
else:
    st.title("🚀 Платная PRO версия🆕")
    
    if not st.session_state.get('authenticated'):
        st.subheader("🔑 Вход в систему")
        license_key = st.text_input("Введите ваш лицензионный ключ", type="password", key="main_license_input")
        
        col_login, col_buy = st.columns(2)
        with col_login:
            if st.button("Войти в аккаунт", use_container_width=True):
                is_active, result = verify_and_update_key(license_key, mode="check")
                if is_active:
                    st.session_state['authenticated'] = True
                    st.session_state['current_user_key'] = license_key
                    st.session_state['remaining'] = result
                    st.success("Успешный вход!")
                    st.rerun()
                else:
                    st.error(result)
        
        with col_buy:
            st.link_button("Купить ключ ($1.9)", "[https://wa.me/992001255656](https://wa.me/992001255656)", use_container_width=True)
            
    else:
        # Панель статуса
        st.success(f"✅ Лицензия активна | Осталось генераций: **{st.session_state.get('remaining', 0)}**")
        
        if st.sidebar.button("Выйти из аккаунта", key="logout_sidebar"):
            st.session_state['authenticated'] = False
            st.rerun()
            
        st.write("---")
        st.subheader("🧠 Генерация через Gemini 2.0 Pro")
        
        pro_topic = st.text_input("О чем будет презентация?", placeholder="Глобальное потепление...", key="pro_topic_field")
        
        c1, c2 = st.columns(2)
        with c1:
            pro_lang = st.selectbox("Язык контента", ["Русский", "English", "Таджикский"], key="pro_lang_sel")
            pro_slides = st.number_input("Количество слайдов", 10, 16, 10, key="pro_slides_val")
        with c2:
            pro_font_size = st.slider("Размер шрифта текста", 20, 30, 22, key="pro_font_slider")
            # Подсказка по лимиту слов
            current_limit = {20:120, 21:110, 22:110, 23:101, 24:95, 25:89, 26:82, 27:77, 28:71, 29:60, 30:55}.get(pro_font_size, 60)
            st.info(f"Слов на слайд: ~{current_limit}")

        st.write("### 🎨 Выберите дизайн")
        templates_dir = "templates"
        style_folders = [f for f in os.listdir(templates_dir) if os.path.isdir(os.path.join(templates_dir, f))]

        if style_folders:
            t_cols = st.columns(3) 
            for idx, folder_name in enumerate(style_folders):
                with t_cols[idx % 3]:
                    folder_path = os.path.join(templates_dir, folder_name)
                    imgs = [f for f in os.listdir(folder_path) if f.endswith(('.png', '.jpg', '.jpeg'))]
                    if imgs:
                        st.image(os.path.join(folder_path, imgs[0]), use_container_width=True)
                    
                    if st.button(folder_name.upper(), key=f"pro_style_{idx}", use_container_width=True):
                        st.session_state['selected_tpl_folder'] = folder_path 
                        st.session_state['pro_style_display'] = folder_name.upper()
            
            if 'pro_style_display' in st.session_state:
                st.markdown(f"<p style='color: #00ff7f;'>✅ Выбран стиль: <b>{st.session_state['pro_style_display']}</b></p>", unsafe_allow_html=True)

        if st.button("🚀 СОЗДАТЬ ПРЕЗЕНТАЦИЮ (1 КРЕДИТ)", key="final_pro_gen_btn", use_container_width=True):
            if not pro_topic:
                st.warning("Сначала введите тему!")
            elif 'selected_tpl_folder' not in st.session_state:
                st.warning("Пожалуйста, выберите дизайн выше!")
            else:
                # 1. Сначала пробуем списать кредит
                user_key = st.session_state.get('current_user_key')
                success_deduct, deduct_res = verify_and_update_key(user_key, mode="update")
                
                if success_deduct:
                    with st.spinner("🤖 Искусственный интеллект создает слайды..."):
                        try:
                            # Генерируем контент
                            ai_json = get_gemini_content(pro_topic, pro_lang, pro_slides, pro_font_size)
                            
                            # Создаем PPTX
                            final_ppt = create_pptx(
                                ai_json,
                                template_folder=st.session_state['selected_tpl_folder'],
                                user_font_size=pro_font_size,
                                topic=pro_topic,
                                include_images=True,
                                image_slide_word_ratio=0.9
                            )
                            
                            if final_ppt:
                                st.session_state['remaining'] = deduct_res
                                st.balloons()
                                st.success(f"Презентация готова! Осталось лимитов: {deduct_res}")
                                
                                st.download_button(
                                    label="📥 СКАЧАТЬ ГОТОВЫЙ ФАЙЛ",
                                    data=final_ppt,
                                    file_name=f"Presentation_{pro_topic}.pptx",
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                    use_container_width=True
                                )
                            else:
                                st.error("Ошибка: ИИ вернул пустой или некорректный ответ. Попробуйте другую тему.")
                                
                        except Exception as e:
                            st.error(f"Ошибка при создании: {e}")
                else:
                    st.error(f"Ошибка списания: {deduct_res}")

st.sidebar.markdown("---")
st.sidebar.markdown(f"**ОСНОВАТЕЛИ: {APP_FOUNDER}**")
